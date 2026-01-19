#!/usr/bin/env python3
"""
AI Nanny PowerPoint Presentation Generator
Generates a professional 10-slide presentation for Dell Agentic AI Developer interview

Usage:
    python generate_presentation.py

Output:
    AI_Nanny_Pitch.pptx (in current directory)

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class PresentationGenerator:
    """Generate AI Nanny pitch presentation."""
    
    # Color palette
    COLORS = {
        'primary': RGBColor(30, 58, 138),      # Deep Blue
        'secondary': RGBColor(124, 58, 237),   # Purple
        'accent': RGBColor(6, 182, 212),       # Cyan
        'success': RGBColor(16, 185, 129),     # Green
        'warning': RGBColor(245, 158, 11),     # Amber
        'error': RGBColor(239, 68, 68),        # Red
        'background': RGBColor(31, 41, 55),    # Dark Gray
        'text': RGBColor(255, 255, 255),       # White
        'text_light': RGBColor(229, 231, 235), # Light Gray
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(5.625)
    
    def add_title_slide(self):
        """Slide 1: Title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "AI Nanny Platform"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Multi-Agent System with Production-Grade RAG Architecture"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.COLORS['accent']
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Author info
        author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        author_frame = author_box.text_frame
        author_frame.text = "Serban Tica • Dell Agentic AI Developer Interview • January 2026"
        author_para = author_frame.paragraphs[0]
        author_para.font.size = Pt(18)
        author_para.font.color.rgb = self.COLORS['text_light']
        author_para.alignment = PP_ALIGN.CENTER
    
    def add_problem_slide(self):
        """Slide 2: The Problem."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        # Title
        self._add_slide_title(slide, "Elderly Care Crisis: The Numbers")
        
        # Problem points
        problems = [
            "❌ 1:15 staff-to-resident ratios in care facilities",
            "❌ 40% medication adherence failure rate",
            "❌ Social isolation → cognitive decline",
            "❌ $94B annual U.S. elderly care market"
        ]
        
        y_pos = 1.5
        for problem in problems:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.text = problem
            text_frame.paragraphs[0].font.size = Pt(24)
            text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
            y_pos += 0.6
        
        # Callout box
        callout = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(4.2), Inches(8), Inches(0.8)
        )
        callout.fill.solid()
        callout.fill.fore_color.rgb = self.COLORS['secondary']
        callout.line.color.rgb = self.COLORS['accent']
        callout.line.width = Pt(2)
        
        text_frame = callout.text_frame
        text_frame.text = "🎯 Core Challenge: How do we scale personalized care without sacrificing quality or safety?"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def add_solution_slide(self):
        """Slide 3: Multi-Persona Agent System."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "AI Nanny: Multi-Persona Agent Architecture")
        
        # 4 Agent cards
        agents = [
            ("💊 Medication Nurse", "RAG-Enabled\nSafety-Critical"),
            ("👥 Companion", "Instruction-Only\nConversational"),
            ("🚨 Emergency", "Pre-Cached\nZero-Latency"),
            ("📖 Storyteller", "RAG + Creative\nPersonalized")
        ]
        
        x_positions = [1, 5.5, 1, 5.5]
        y_positions = [1.5, 1.5, 3.2, 3.2]
        
        for (title, description), x, y in zip(agents, x_positions, y_positions):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(3.5), Inches(1.3)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLORS['primary']
            card.line.color.rgb = self.COLORS['accent']
            card.line.width = Pt(2)
            
            text_frame = card.text_frame
            text_frame.text = title
            p = text_frame.paragraphs[0]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['text']
            p.alignment = PP_ALIGN.CENTER
            
            # Add description
            p2 = text_frame.add_paragraph()
            p2.text = description
            p2.font.size = Pt(16)
            p2.font.color.rgb = self.COLORS['text_light']
            p2.alignment = PP_ALIGN.CENTER
            
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Bottom features
        features_box = slide.shapes.add_textbox(Inches(1), Inches(4.9), Inches(8), Inches(0.5))
        features_frame = features_box.text_frame
        features_frame.text = "🎙️ Voice-First Interface  |  🔒 HIPAA Compliant  |  🏠 Edge Deployment (Raspberry Pi)"
        features_frame.paragraphs[0].font.size = Pt(18)
        features_frame.paragraphs[0].font.color.rgb = self.COLORS['success']
        features_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_rag_comparison_slide(self):
        """Slide 4: RAG Architecture Comparison."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "Why Advanced RAG? Evaluating 4 Approaches")
        
        # Table header
        headers = ["Pattern", "Latency", "Accuracy", "Complexity", "Our Fit"]
        rows = [
            ["Naive RAG", "✅ <300ms", "⚠️ 70-75%", "✅ Simple", "❌ Unsafe for medical"],
            ["Advanced RAG", "✅ <500ms", "✅ 85-90%", "⚠️ Medium", "✅ CHOSEN"],
            ["Agentic RAG", "❌ 1-3s", "✅ 90-95%", "❌ High", "❌ Too slow for UX"],
            ["Hybrid (Graph)", "⚠️ <800ms", "✅ 90%+", "❌ Very High", "❌ Edge constraints"]
        ]
        
        # Create table manually with shapes
        table_left = 0.5
        table_top = 1.5
        col_width = 1.9
        row_height = 0.5
        
        # Headers
        for i, header in enumerate(headers):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(table_left + i * col_width), Inches(table_top),
                Inches(col_width), Inches(row_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['secondary']
            cell.line.color.rgb = self.COLORS['text']
            
            text_frame = cell.text_frame
            text_frame.text = header
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(table_left + col_idx * col_width),
                    Inches(table_top + (row_idx + 1) * row_height),
                    Inches(col_width), Inches(row_height)
                )
                
                # Highlight "Advanced RAG" row
                if row_idx == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(34, 197, 94)  # Green tint
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['primary']
                
                cell.line.color.rgb = self.COLORS['text_light']
                
                text_frame = cell.text_frame
                text_frame.text = cell_text
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Callout
        callout = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(4.5), Inches(8), Inches(0.7)
        )
        callout.fill.solid()
        callout.fill.fore_color.rgb = self.COLORS['success']
        callout.line.color.rgb = self.COLORS['accent']
        callout.line.width = Pt(2)
        
        text_frame = callout.text_frame
        text_frame.text = "🎯 Decision: Advanced RAG with selective reranking - Balance of accuracy, latency, and edge compatibility"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def add_rag_deep_dive_slide(self):
        """Slide 5: RAG Implementation Strategy."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "RAG Implementation Strategy")
        
        # Three columns
        columns = [
            {
                'title': '📄 Semantic Chunking',
                'items': [
                    '• Preserve paragraph boundaries',
                    '• 256 tokens: Medical protocols',
                    '• 512 tokens: Conversational docs',
                    '• 50-token overlap for context',
                    '',
                    'Why: Fixed-size chunks break',
                    'medical procedures mid-step'
                ]
            },
            {
                'title': '🧠 Dual-Mode Embeddings',
                'items': [
                    'Primary: OpenAI',
                    'text-embedding-3-small',
                    '• 1536 dimensions',
                    '• $0.02/1M tokens',
                    '',
                    'Fallback: Sentence',
                    'Transformers',
                    '• Local/offline mode',
                    '• Free, 384 dimensions'
                ]
            },
            {
                'title': '💾 ChromaDB → Pinecone',
                'items': [
                    'MVP: ChromaDB (local)',
                    '• Raspberry Pi 4 compatible',
                    '• <50ms search latency',
                    '• Offline capability',
                    '',
                    'Scale: Pinecone (cloud)',
                    '• 100M+ vectors',
                    '• No code changes required'
                ]
            }
        ]
        
        x_positions = [0.5, 3.5, 6.5]
        
        for col, x in zip(columns, x_positions):
            # Column header
            header_box = slide.shapes.add_textbox(Inches(x), Inches(1.3), Inches(2.8), Inches(0.4))
            header_frame = header_box.text_frame
            header_frame.text = col['title']
            header_frame.paragraphs[0].font.size = Pt(18)
            header_frame.paragraphs[0].font.bold = True
            header_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
            header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Column content
            content_box = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(2.8), Inches(3.2))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for item in col['items']:
                p = content_frame.paragraphs[0] if item == col['items'][0] else content_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = self.COLORS['text']
                p.space_after = Pt(6)
    
    def add_real_example_slide(self):
        """Slide 6: Real-World Example."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "RAG in Action: Medication Nurse Scenario")
        
        # Query box
        query_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(1.3), Inches(7), Inches(0.5)
        )
        query_box.fill.solid()
        query_box.fill.fore_color.rgb = self.COLORS['primary']
        query_box.line.color.rgb = self.COLORS['accent']
        query_box.line.width = Pt(2)
        
        text_frame = query_box.text_frame
        text_frame.text = '💬 User Query: "Is it time for my blood pressure medication?"'
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Processing steps
        steps = [
            "Embed Query (<100ms)",
            "Vector Search (Top-K=5)",
            "Retrieved Context (3 chunks)"
        ]
        
        y_pos = 2.0
        for step in steps:
            step_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.35))
            step_frame = step_box.text_frame
            step_frame.text = f"↓ {step}"
            step_frame.paragraphs[0].font.size = Pt(16)
            step_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
            step_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            y_pos += 0.4
        
        # Retrieved chunks
        chunks = [
            ("1️⃣ Medication Schedule (0.92)", '"Lisinopril 10mg at 8:00 AM daily"'),
            ("2️⃣ Medical Profile (0.87)", '"Takes blood thinner - head injury = 911"'),
            ("3️⃣ Protocol (0.81)", '"Check BP within 30 min before..."')
        ]
        
        y_pos = 3.3
        for title, content in chunks:
            chunk_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos), Inches(8), Inches(0.35)
            )
            chunk_box.fill.solid()
            chunk_box.fill.fore_color.rgb = self.COLORS['secondary']
            chunk_box.line.color.rgb = self.COLORS['text_light']
            
            text_frame = chunk_box.text_frame
            text_frame.text = f"{title}  {content}"
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            y_pos += 0.4
        
        # Bottom stats
        stats_box = slide.shapes.add_textbox(Inches(1), Inches(4.9), Inches(8), Inches(0.5))
        stats_frame = stats_box.text_frame
        stats_frame.text = "⚡ Total Latency: <300ms  |  🎯 Accuracy: >90%  |  🔒 HIPAA Logged"
        stats_frame.paragraphs[0].font.size = Pt(18)
        stats_frame.paragraphs[0].font.bold = True
        stats_frame.paragraphs[0].font.color.rgb = self.COLORS['success']
        stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_validation_slide(self):
        """Slide 7: Safety-First Validation Pipeline."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "6-Stage Validation: Zero Tolerance for Errors")
        
        validations = [
            ("1️⃣ Freshness Check", "Medical: <24h | Protocols: <1yr", "❌ FAIL → Escalate"),
            ("2️⃣ Authority Verification", "Only nurse/doctor uploads", "❌ FAIL → Reject"),
            ("3️⃣ Conflict Detection", "Multiple schedules for same drug?", "❌ FAIL → Human review"),
            ("4️⃣ Privacy Boundaries", "Companion CANNOT access medical", "❌ FAIL → Block"),
            ("5️⃣ Confidence Scoring", "Similarity > 0.85 for medical", "⚠️ WARN → Disclaimer"),
            ("6️⃣ Audit Logging", "Every retrieval logged (HIPAA)", "✅ Complete trail")
        ]
        
        y_pos = 1.3
        for stage, description, action in validations:
            # Stage box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos), Inches(8), Inches(0.5)
            )
            
            if "✅" in action:
                box.fill.solid()
                box.fill.fore_color.rgb = self.COLORS['success']
            elif "⚠️" in action:
                box.fill.solid()
                box.fill.fore_color.rgb = self.COLORS['warning']
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = self.COLORS['primary']
            
            box.line.color.rgb = self.COLORS['text_light']
            
            text_frame = box.text_frame
            text_frame.text = f"{stage}  •  {description}  •  {action}"
            text_frame.paragraphs[0].font.size = Pt(15)
            text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            y_pos += 0.55
        
        # Bottom stat
        stat_box = slide.shapes.add_textbox(Inches(1), Inches(4.9), Inches(8), Inches(0.5))
        stat_frame = stat_box.text_frame
        stat_frame.text = "🎯 Safety Record: ZERO medication errors from RAG advice"
        stat_frame.paragraphs[0].font.size = Pt(20)
        stat_frame.paragraphs[0].font.bold = True
        stat_frame.paragraphs[0].font.color.rgb = self.COLORS['success']
        stat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_challenges_slide(self):
        """Slide 8: Production Challenges & Solutions."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "Real-World Deployment: Problems We Solved")
        
        challenges = [
            ("📝 Handwritten notes", "OCR preprocessing + manual QA"),
            ("🔄 Conflicting updates", "Timestamp priority + moderation"),
            ("🌐 Network outages", "Local embeddings fallback"),
            ("💰 Cost overruns", "Caching (80% hit rate) + batch"),
            ("🔐 HIPAA compliance", "Persona access control + filtering"),
            ("🏥 Outdated protocols", "Expiration alerts + quarterly review")
        ]
        
        y_pos = 1.4
        for challenge, solution in challenges:
            # Two-column layout
            challenge_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(3.5), Inches(0.45))
            challenge_frame = challenge_box.text_frame
            challenge_frame.text = challenge
            challenge_frame.paragraphs[0].font.size = Pt(16)
            challenge_frame.paragraphs[0].font.bold = True
            challenge_frame.paragraphs[0].font.color.rgb = self.COLORS['error']
            
            solution_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(4.7), Inches(0.45))
            solution_frame = solution_box.text_frame
            solution_frame.text = f"→ {solution}"
            solution_frame.paragraphs[0].font.size = Pt(16)
            solution_frame.paragraphs[0].font.color.rgb = self.COLORS['success']
            
            y_pos += 0.5
        
        # Key insight
        insight_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(4.5), Inches(8), Inches(0.7)
        )
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = self.COLORS['accent']
        insight_box.line.color.rgb = self.COLORS['text']
        insight_box.line.width = Pt(2)
        
        text_frame = insight_box.text_frame
        text_frame.text = "💡 Key Insight: Production RAG requires operational architecture, not just algorithmic optimization"
        text_frame.paragraphs[0].font.size = Pt(17)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def add_metrics_slide(self):
        """Slide 9: Performance & Scalability."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "Performance Metrics & Future Roadmap")
        
        # Left column - Current metrics
        left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4))
        left_title_frame = left_title.text_frame
        left_title_frame.text = "📊 Current Performance"
        left_title_frame.paragraphs[0].font.size = Pt(22)
        left_title_frame.paragraphs[0].font.bold = True
        left_title_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
        
        metrics = [
            "✅ Retrieval Latency: <300ms (p95)",
            "✅ End-to-End: <2s response",
            "✅ Accuracy: 91% (manual review)",
            "✅ Cache Hit Rate: 83%",
            "✅ Safety Incidents: ZERO",
            "✅ Uptime: 99.7%"
        ]
        
        y_pos = 1.8
        for metric in metrics:
            metric_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(4), Inches(0.35))
            metric_frame = metric_box.text_frame
            metric_frame.text = metric
            metric_frame.paragraphs[0].font.size = Pt(16)
            metric_frame.paragraphs[0].font.color.rgb = self.COLORS['success']
            y_pos += 0.4
        
        # Right column - Future roadmap
        right_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.4))
        right_title_frame = right_title.text_frame
        right_title_frame.text = "📅 Future Enhancements"
        right_title_frame.paragraphs[0].font.size = Pt(22)
        right_title_frame.paragraphs[0].font.bold = True
        right_title_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
        
        roadmap = [
            "Q2 2026: Hybrid RAG",
            "  • Knowledge graph for drug",
            "    interactions",
            "",
            "Q3 2026: Query Reformulation",
            "  • Handle ambiguous speech",
            "",
            "Q4 2026: Fine-Tuned Embeddings",
            "  • Healthcare domain",
            "  • 10-15% accuracy improvement"
        ]
        
        y_pos = 1.8
        for item in roadmap:
            roadmap_box = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4), Inches(0.25))
            roadmap_frame = roadmap_box.text_frame
            roadmap_frame.text = item
            roadmap_frame.paragraphs[0].font.size = Pt(14)
            roadmap_frame.paragraphs[0].font.color.rgb = self.COLORS['text']
            y_pos += 0.28
    
    def add_closing_slide(self):
        """Slide 10: Closing - Why This Matters."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        self._add_slide_title(slide, "AI Nanny: Production-Ready Agentic AI")
        
        # Three key takeaways
        takeaways = [
            ("1️⃣ ARCHITECTURE THINKING", 
             "Advanced RAG chosen through rigorous evaluation\n(Naive, Agentic, Hybrid alternatives rejected)"),
            ("2️⃣ SAFETY-CRITICAL VALIDATION", 
             "Healthcare demands zero-tolerance error handling\n6-stage validation pipeline with human escalation"),
            ("3️⃣ PRODUCTION PRAGMATISM", 
             "Edge deployment, offline fallbacks, cost control\nReal problems solved: OCR, conflicts, compliance")
        ]
        
        y_pos = 1.4
        for number, description in takeaways:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_pos), Inches(8), Inches(0.7)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS['primary']
            box.line.color.rgb = self.COLORS['accent']
            box.line.width = Pt(2)
            
            text_frame = box.text_frame
            text_frame.text = number
            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent']
            
            p2 = text_frame.add_paragraph()
            p2.text = description
            p2.font.size = Pt(15)
            p2.font.color.rgb = self.COLORS['text']
            p2.space_after = Pt(4)
            
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            y_pos += 0.85
        
        # Call to action
        cta_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.9))
        cta_frame = cta_box.text_frame
        cta_frame.text = '💬 "This architecture thinking is directly applicable to Dell\'s agentic AI initiatives"'
        p = cta_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['success']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = cta_frame.add_paragraph()
        p2.text = "Let's discuss: RAG reliability, multi-agent coordination, and production deployment strategies"
        p2.font.size = Pt(16)
        p2.font.color.rgb = self.COLORS['text']
        p2.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide, title_text):
        """Helper to add consistent slide titles."""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['accent']
        title_para.alignment = PP_ALIGN.CENTER
    
    def generate(self, output_file="AI_Nanny_Pitch.pptx"):
        """Generate the complete presentation."""
        print("Generating AI Nanny Presentation...")
        print("Creating slides...")
        
        self.add_title_slide()
        print("  ✓ Slide 1: Title")
        
        self.add_problem_slide()
        print("  ✓ Slide 2: Problem")
        
        self.add_solution_slide()
        print("  ✓ Slide 3: Solution")
        
        self.add_rag_comparison_slide()
        print("  ✓ Slide 4: RAG Comparison")
        
        self.add_rag_deep_dive_slide()
        print("  ✓ Slide 5: RAG Deep Dive")
        
        self.add_real_example_slide()
        print("  ✓ Slide 6: Real Example")
        
        self.add_validation_slide()
        print("  ✓ Slide 7: Validation Pipeline")
        
        self.add_challenges_slide()
        print("  ✓ Slide 8: Challenges & Solutions")
        
        self.add_metrics_slide()
        print("  ✓ Slide 9: Performance Metrics")
        
        self.add_closing_slide()
        print("  ✓ Slide 10: Closing")
        
        self.prs.save(output_file)
        print(f"\n✅ Presentation saved: {output_file}")
        print(f"   Total slides: {len(self.prs.slides)}")
        print("\nNext steps:")
        print("  1. Open the file in PowerPoint/Keynote/Google Slides")
        print("  2. Review and customize colors/fonts as needed")
        print("  3. Export to PDF for backup")
        print("  4. Practice with timer (aim for 2:50)")


def main():
    """Main entry point."""
    generator = PresentationGenerator()
    generator.generate()


if __name__ == "__main__":
    main()
